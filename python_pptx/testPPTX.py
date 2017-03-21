# -*- coding: utf-8 -*-
#!/usr/bin/python
#!/usr/bin/env python

#pip install python-pptx

from pptx import Presentation
from pptx.util import Inches
import sys
import os

def getImageList(source):
    matches = []
    for root, dirnames, filenames in os.walk(source):
        for filename in filenames:
            if filename.endswith(('.png', '.PNG', '.jpg', '.JPG','.jpeg', '.JPEG', '.gif', '.GIF')):
                matches.append(os.path.join(root, filename))
    return matches

def main():
    ## INPUT PARAMETERS
    outputPath    = 'output/test.pptx'
    inputTemplate = 'template/template.pptx'
    prs = Presentation(inputTemplate)


    bullet_slide_layout = prs.slide_layouts[1] # La slide de layout debe ser distinta a la slide de inicio de presentacion
    left = Inches(1)
    top  = Inches(1)

    newImgPath = "E:\sandbox-python\python_pptx\imgs"

    allImgsNames = getImageList(newImgPath)
    print allImgsNames

    allImgsNamesLen = len(allImgsNames)

    for i in range(0, allImgsNamesLen):
        newImgPath = allImgsNames[i]
        newImgPath = newImgPath.replace('\\','/')
        newImgPathList = newImgPath.split('/')
        newImgName = newImgPathList[-1]
        newImgNameList = newImgName.split('.')
        newImgNameList = newImgNameList[:-1]
        sep = "."
        newImgName     = sep.join(newImgNameList)
        print "Imagen " + str(i) + " de " + str(allImgsNamesLen) + ": " + newImgName

         # Se agrega una slide en base a layout
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape      = shapes.title
        body_shape       = shapes.placeholders[1]
        title_shape.text = newImgName
        ## Se inserta la imagen
        try:
            pic = slide.shapes.add_picture(newImgPath, left, top)
        except IOError:
            print "ERROR: El archivo " + newImgPath + " no existe"
            sys.exit()
        ## Centrar imagen de entrada en la slide
        pic.left         = (prs.slide_width - pic.width) / 2
        pic.top          = (prs.slide_height - pic.height) / 2

    ## SAVE FILE
    print "Guardando archivo " + outputPath
    prs.save(outputPath)
    print "DONE!"

if __name__ == "__main__":
    main()